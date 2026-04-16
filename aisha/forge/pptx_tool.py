"""PPTX generation + Slack upload.

Minimal surface: generate a deck from a structured outline (title + slides of
bullets), optionally upload it to a Slack channel/thread via the Slack
``files.upload_v2`` flow.

One module, stdlib HTTP, no extra deps beyond python-pptx.
"""
from __future__ import annotations

import json
import logging
import os
import tempfile
import urllib.parse
import urllib.request
import uuid
from pathlib import Path
from typing import Optional

log = logging.getLogger(__name__)

_SLACK_API = "https://slack.com/api"


# ── Deck generation ────────────────────────────────────────────────────

def generate_pptx(
    title: str,
    slides: list[dict],
    *,
    subtitle: str = "",
    out_path: Optional[Path] = None,
) -> Path:
    """Build a .pptx from a structured outline.

    ``slides`` is a list of ``{"title": str, "bullets": [str, ...]}`` dicts.
    Returns the path to the generated file.
    """
    from pptx import Presentation  # lazy import — keeps chat boot fast

    prs = Presentation()

    # Title slide
    title_layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(title_layout)
    s.shapes.title.text = title
    if subtitle and len(s.placeholders) > 1:
        s.placeholders[1].text = subtitle

    # Content slides — one per entry
    content_layout = prs.slide_layouts[1]  # "Title and Content"
    for item in slides:
        s = prs.slides.add_slide(content_layout)
        s.shapes.title.text = item.get("title", "").strip() or "Slide"
        body = s.placeholders[1].text_frame
        body.clear()
        bullets = item.get("bullets") or []
        if not bullets:
            bullets = [""]
        for i, b in enumerate(bullets):
            text = str(b).strip()
            if i == 0:
                body.paragraphs[0].text = text
            else:
                p = body.add_paragraph()
                p.text = text
                p.level = 0

    if out_path is None:
        name = f"aisha_{uuid.uuid4().hex[:8]}.pptx"
        out_path = Path(tempfile.gettempdir()) / name
    prs.save(str(out_path))
    log.info("pptx: generated %s (%d slides)", out_path, len(slides))
    return out_path


# ── Slack upload (files.upload_v2 flow) ────────────────────────────────

def _slack_post(path: str, token: str, body: dict, *, json_body: bool = False) -> dict:
    url = f"{_SLACK_API}/{path}"
    headers = {"Authorization": f"Bearer {token}"}
    if json_body:
        data = json.dumps(body).encode("utf-8")
        headers["Content-Type"] = "application/json; charset=utf-8"
    else:
        data = urllib.parse.urlencode(body).encode("utf-8")
        headers["Content-Type"] = "application/x-www-form-urlencoded"
    req = urllib.request.Request(url, data=data, headers=headers, method="POST")
    with urllib.request.urlopen(req, timeout=30) as resp:
        return json.loads(resp.read().decode("utf-8"))


def upload_to_slack(
    path: Path,
    *,
    channel: str,
    thread_ts: Optional[str] = None,
    title: Optional[str] = None,
    initial_comment: Optional[str] = None,
) -> dict:
    """Upload a file to a Slack channel/thread via files.upload_v2.

    Returns ``{"ok": bool, "permalink": str, ...}`` on success; raises on
    network errors. Logs and returns ``{"ok": False, ...}`` on API errors.
    """
    token = os.environ.get("SLACK_BOT_TOKEN", "")
    if not token:
        return {"ok": False, "error": "SLACK_BOT_TOKEN not set"}

    path = Path(path)
    size = path.stat().st_size
    filename = path.name

    # Step 1: get a presigned upload URL + file id
    step1 = _slack_post(
        "files.getUploadURLExternal",
        token,
        {"filename": filename, "length": str(size)},
    )
    if not step1.get("ok"):
        log.error("slack upload step1 failed: %s", step1)
        return {"ok": False, "error": step1.get("error", "getUploadURLExternal failed")}
    upload_url = step1["upload_url"]
    file_id = step1["file_id"]

    # Step 2: POST the raw bytes to the presigned URL (no auth header).
    with path.open("rb") as f:
        payload = f.read()
    req = urllib.request.Request(
        upload_url,
        data=payload,
        headers={"Content-Type": "application/octet-stream"},
        method="POST",
    )
    try:
        with urllib.request.urlopen(req, timeout=120) as resp:
            resp.read()
    except Exception as e:
        log.error("slack upload step2 (binary POST) failed: %s", e)
        return {"ok": False, "error": f"upload POST failed: {e}"}

    # Step 3: complete + share to channel/thread
    body: dict = {
        "files": [{"id": file_id, "title": title or filename}],
        "channel_id": channel,
    }
    if thread_ts:
        body["thread_ts"] = thread_ts
    if initial_comment:
        body["initial_comment"] = initial_comment
    step3 = _slack_post("files.completeUploadExternal", token, body, json_body=True)
    if not step3.get("ok"):
        log.error("slack upload step3 failed: %s", step3)
        return {"ok": False, "error": step3.get("error", "completeUploadExternal failed")}

    files = step3.get("files") or [{}]
    permalink = files[0].get("permalink", "")
    log.info("pptx: uploaded to slack channel=%s thread=%s permalink=%s",
             channel, thread_ts, permalink)
    return {"ok": True, "permalink": permalink, "file_id": file_id}


# ── Generic Slack post ─────────────────────────────────────────────────

def update_message(channel: str, ts: str, text: str) -> dict:
    """Edit an existing message via chat.update. Bot can only edit its own messages."""
    token = os.environ.get("SLACK_BOT_TOKEN", "")
    if not token:
        return {"ok": False, "error": "SLACK_BOT_TOKEN not set"}
    if not channel or not ts or not text:
        return {"ok": False, "error": "channel, ts, and text required"}
    try:
        resp = _slack_post(
            "chat.update", token,
            {"channel": channel, "ts": ts, "text": text},
            json_body=True,
        )
    except Exception as e:
        log.error("update_message failed: %s", e)
        return {"ok": False, "error": str(e)}
    if not resp.get("ok"):
        log.warning("update_message: slack returned error: %s", resp)
    return resp


def delete_message(channel: str, ts: str) -> dict:
    """Delete an existing message via chat.delete. Bot can only delete its own."""
    token = os.environ.get("SLACK_BOT_TOKEN", "")
    if not token:
        return {"ok": False, "error": "SLACK_BOT_TOKEN not set"}
    if not channel or not ts:
        return {"ok": False, "error": "channel and ts required"}
    try:
        resp = _slack_post(
            "chat.delete", token,
            {"channel": channel, "ts": ts},
            json_body=True,
        )
    except Exception as e:
        log.error("delete_message failed: %s", e)
        return {"ok": False, "error": str(e)}
    if not resp.get("ok"):
        log.warning("delete_message: slack returned error: %s", resp)
    return resp


def post_message(
    channel: str,
    text: str,
    *,
    thread_ts: Optional[str] = None,
) -> dict:
    """Send a message via chat.postMessage. Channel may be an id (C…) or a name.

    Returns ``{"ok": bool, "ts": str, "channel": str, "error": str}``. Slack
    resolves channel names server-side when prefixed with ``#`` — pass the
    literal string the user provided and let Slack decide.
    """
    token = os.environ.get("SLACK_BOT_TOKEN", "")
    if not token:
        return {"ok": False, "error": "SLACK_BOT_TOKEN not set"}
    if not channel or not text:
        return {"ok": False, "error": "channel and text required"}
    body: dict = {"channel": channel, "text": text}
    if thread_ts:
        body["thread_ts"] = thread_ts
    try:
        resp = _slack_post("chat.postMessage", token, body, json_body=True)
    except Exception as e:
        log.error("post_message failed: %s", e)
        return {"ok": False, "error": str(e)}
    if not resp.get("ok"):
        log.warning("post_message: slack returned error: %s", resp)
    return resp


# ── Source-tag parsing ─────────────────────────────────────────────────

def slack_context(source: str) -> Optional[dict]:
    """Parse a ``slack:CHANNEL:THREAD_TS`` source tag into dict parts.

    Returns None if ``source`` isn't a Slack tag.
    """
    if not source or not source.startswith("slack:"):
        return None
    parts = source.split(":", 2)
    if len(parts) < 2:
        return None
    channel = parts[1]
    thread_ts = parts[2] if len(parts) > 2 and parts[2] else None
    return {"channel": channel, "thread_ts": thread_ts}
